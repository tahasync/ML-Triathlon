import os
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

BASE = r'C:\Users\Tahan\Desktop\Claude\AI Project'
RESULTS = os.path.join(BASE, 'results')

# ═══════════════════════════════════════════
#  REPORT DATA
# ═══════════════════════════════════════════

REPORT_DATA = {
    'kmeans': {
        'title': 'K-Means Clustering on LEGO Database',
        'dataset': 'LEGO Database (Rebrickable)',
        'dataset_desc': '26,989 LEGO sets with attributes: set number, name, year, theme, number of parts',
        'algorithm': 'K-Means Clustering',
        'objective': 'Cluster LEGO sets by year and number of parts to discover hidden patterns in LEGO hobbyist trends.',
        'features': 'year (release year), num_parts (piece count)',
        'preprocessing': 'StandardScaler normalization; filtered to sets from 1950 onward with >0 parts',
        'optimal_k': 4,
        'silhouette': 0.5202,
        'cluster_desc': [
            ('Cluster 0 — Modern Small Sets', '12,104 sets', 'Modern small/accessory sets (polybags, books, electronics like NXT bricks). Low part count, recent years.'),
            ('Cluster 1 — Modern Large Collector Sets', '201 sets', 'High piece count (2000+) modern sets like the Burrow Collectors\' Edition, large Ninjago/Icons builds. Premium/collector segment.'),
            ('Cluster 2 — Vintage Sets', '5,677 sets', 'Classic sets from 1950s-1980s with very low part counts. Includes early Technic, Samsonite baseplates, Playhouse. Represents LEGO\'s early years.'),
            ('Cluster 3 — Modern Mid-Size Sets', '1,830 sets', 'Modern sets in the 500-1500 piece range. Mainstream retail sets from Star Wars, Ninjago, Harry Potter themes.'),
        ],
        'images': ['lego_elbow.png', 'lego_clusters.png'],
        'conclusion': 'K-Means effectively separates LEGO sets into 4 meaningful clusters based on era and complexity. The silhouette score of 0.52 indicates moderately well-separated clusters. The vintage cluster (Cluster 2) and large collector sets (Cluster 1) are the most distinct groups.'
    },
    'knn': {
        'title': 'K-Nearest Neighbors on Aviation Wildlife Strikes',
        'dataset': 'FAA Wildlife Strikes Database (OpenIntro)',
        'dataset_desc': '19,302 wildlife strike records from 1990-1997 with aircraft, species, and environmental details',
        'algorithm': 'K-Nearest Neighbors (KNN) Classification',
        'objective': 'Classify the risk level of a wildlife strike based on aircraft data, species, and environmental conditions.',
        'features': 'height (ft), speed (knots), aircraft mass class, number of engines, phase of flight, time of day, sky condition, species',
        'target': 'Effect on flight (None, Other, Precautionary Landing, Aborted Take-off, Engine Shut Down)',
        'preprocessing': 'Label encoding for categorical features; median imputation for numerical features; StandardScaler normalization',
        'best_k': 8,
        'cv_accuracy': 0.8957,
        'test_accuracy': 0.8977,
        'class_distribution': [
            ('None', '17,329 (89.8%)'),
            ('Other', '348 (1.8%)'),
            ('Precautionary Landing', '965 (5.0%)'),
            ('Aborted Take-off', '544 (2.8%)'),
            ('Engine Shut Down', '116 (0.6%)'),
        ],
        'images': ['wildlife_knn_accuracy.png', 'wildlife_confusion_matrix.png'],
        'conclusion': 'The KNN classifier achieves 89.8% accuracy with k=8. However, the dataset is highly imbalanced — 89.8% of strikes result in no effect. The model excels at identifying non-damaging strikes (recall=1.00) but struggles with rare high-severity classes. Future improvements could include SMOTE oversampling or class-weighted training.'
    },
    'naive_bayes': {
        'title': 'Naive Bayes on UFO Sightings Reports',
        'dataset': 'NUFORC UFO Sightings Database',
        'dataset_desc': '80,331 UFO sighting reports with datetime, location, shape, duration, and text description',
        'algorithm': 'Multinomial Naive Bayes with TF-IDF Vectorization',
        'objective': 'Predict the reported shape of a UFO based on the textual description of the sighting.',
        'features': 'TF-IDF vectors (5,000 features, unigrams + bigrams) from sighting comments',
        'target': 'UFO shape (10 classes: light, triangle, circle, fireball, other, unknown, sphere, disk, oval, formation)',
        'preprocessing': 'Lowercasing, URL removal, punctuation stripping, stopword removal, short word filtering; TF-IDF with sublinear tf',
        'test_accuracy': 0.4610,
        'top_shapes': [
            ('light', '16,565 (24.3%)'),
            ('triangle', '7,865 (11.5%)'),
            ('circle', '7,608 (11.2%)'),
            ('fireball', '6,208 (9.1%)'),
            ('other', '5,649 (8.3%)'),
        ],
        'best_classes': [
            ('Triangle', '0.61 F1-score — distinctive descriptions ("silent", "triangular", "black")'),
            ('Fireball', '0.57 F1-score — strong signal words ("orange", "streak", "explosion")'),
            ('Light', '0.56 F1-score — most frequent class with generic terms'),
            ('Disk', '0.48 F1-score — classic "saucer" descriptions'),
            ('Sphere', '0.45 F1-score — "orb", "glowing ball" patterns'),
        ],
        'images': ['ufo_confusion_matrix.png', 'ufo_top_features.png'],
        'conclusion': 'Multinomial Naive Bayes achieves 46.1% accuracy across 10 shape classes, significantly outperforming the majority-class baseline of ~24%. Triangle, fireball, and light shapes are most reliably identified. The confusion matrix shows expected confusion between visually similar shapes (e.g., circle vs sphere, light vs unknown). This demonstrates that sighting descriptions contain meaningful signal for shape classification.'
    }
}

IMAGES = {
    'lego_elbow.png': os.path.join(RESULTS, 'lego_elbow.png'),
    'lego_clusters.png': os.path.join(RESULTS, 'lego_clusters.png'),
    'wildlife_knn_accuracy.png': os.path.join(RESULTS, 'wildlife_knn_accuracy.png'),
    'wildlife_confusion_matrix.png': os.path.join(RESULTS, 'wildlife_confusion_matrix.png'),
    'ufo_confusion_matrix.png': os.path.join(RESULTS, 'ufo_confusion_matrix.png'),
    'ufo_top_features.png': os.path.join(RESULTS, 'ufo_top_features.png'),
}

# ═══════════════════════════════════════════
#  DOCX REPORT
# ═══════════════════════════════════════════

def add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    return h

def add_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Light Grid Accent 1'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            table.rows[ri + 1].cells[ci].text = str(val)
    return table

def add_image(doc, img_name, width=Inches(5.5)):
    path = IMAGES.get(img_name)
    if path and os.path.exists(path):
        doc.add_picture(path, width=width)
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        caption = doc.add_paragraph(img_name.replace('.png', '').replace('_', ' ').title())
        caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in caption.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

def build_report():
    doc = Document()

    # Title page
    title = doc.add_heading('AI Model Training Report', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = doc.add_paragraph('Comparative Analysis of K-Means, KNN, and Naive Bayes Algorithms')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in subtitle.runs:
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    doc.add_paragraph('')
    meta = doc.add_paragraph('Course: Artificial Intelligence\nDate: June 2026')
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # Table of Contents
    add_heading(doc, 'Table of Contents', level=1)
    toc_items = [
        '1. Introduction',
        '2. K-Means Clustering — LEGO Database',
        '3. K-Nearest Neighbors — Aviation Wildlife Strikes',
        '4. Naive Bayes — UFO Sightings Reports',
        '5. Comparative Analysis',
        '6. Conclusion',
    ]
    for item in toc_items:
        doc.add_paragraph(item, style='List Number')

    doc.add_page_break()

    # 1. Introduction
    add_heading(doc, '1. Introduction', level=1)
    doc.add_paragraph(
        'This report presents the implementation and evaluation of three machine learning algorithms '
        'applied to three distinct real-world datasets. Each algorithm was selected based on its suitability '
        'for the specific problem domain and dataset characteristics.'
    )
    add_table(doc,
        ['Algorithm', 'Dataset', 'Problem Type', 'Objective'],
        [
            ['K-Means', 'LEGO Database', 'Unsupervised Clustering', 'Group LEGO sets by complexity and era'],
            ['KNN', 'Aviation Wildlife Strikes', 'Supervised Classification', 'Classify strike risk level'],
            ['Naive Bayes', 'UFO Sightings Reports', 'Text Classification', 'Predict UFO shape from descriptions'],
        ]
    )
    doc.add_page_break()

    # 2. K-Means
    d = REPORT_DATA['kmeans']
    add_heading(doc, f'2. {d["title"]}', level=1)
    doc.add_paragraph(f'Dataset: {d["dataset"]}')
    doc.add_paragraph(d['dataset_desc'])
    doc.add_paragraph(f'Algorithm: {d["algorithm"]}')
    doc.add_paragraph(f'Objective: {d["objective"]}')

    add_heading(doc, '2.1 Methodology', level=2)
    doc.add_paragraph(f'Features used: {d["features"]}')
    doc.add_paragraph(f'Preprocessing: {d["preprocessing"]}')
    doc.add_paragraph(
        'The optimal number of clusters was determined using the Elbow Method (inertia) '
        'and Silhouette Score analysis for k values from 2 to 10.'
    )

    add_heading(doc, '2.2 Results', level=2)
    doc.add_paragraph(f'Optimal k: {d["optimal_k"]}')
    doc.add_paragraph(f'Silhouette Score: {d["silhouette"]}')

    add_heading(doc, '2.3 Cluster Analysis', level=3)
    for name, size, desc in d['cluster_desc']:
        p = doc.add_paragraph()
        run = p.add_run(f'{name} ({size})')
        run.bold = True
        p.add_run(f': {desc}')

    add_heading(doc, '2.4 Visualizations', level=2)
    for img in d['images']:
        add_image(doc, img)

    add_heading(doc, '2.5 Conclusion', level=2)
    doc.add_paragraph(d['conclusion'])
    doc.add_page_break()

    # 3. KNN
    d = REPORT_DATA['knn']
    add_heading(doc, f'3. {d["title"]}', level=1)
    doc.add_paragraph(f'Dataset: {d["dataset"]}')
    doc.add_paragraph(d['dataset_desc'])
    doc.add_paragraph(f'Algorithm: {d["algorithm"]}')
    doc.add_paragraph(f'Objective: {d["objective"]}')

    add_heading(doc, '3.1 Methodology', level=2)
    doc.add_paragraph(f'Features: {d["features"]}')
    doc.add_paragraph(f'Target: {d["target"]}')
    doc.add_paragraph(f'Preprocessing: {d["preprocessing"]}')
    doc.add_paragraph('Optimal k was selected by evaluating accuracy on a held-out test set for k = 1 to 30.')

    add_heading(doc, '3.2 Results', level=2)
    doc.add_paragraph(f'Best k: {d["best_k"]}')
    doc.add_paragraph(f'Test Accuracy: {d["test_accuracy"]}')
    doc.add_paragraph(f'5-Fold Cross-Validation Accuracy: {d["cv_accuracy"]}')

    add_heading(doc, '3.3 Class Distribution', level=3)
    add_table(doc, ['Class', 'Count'], d['class_distribution'])

    add_heading(doc, '3.4 Visualizations', level=2)
    for img in d['images']:
        add_image(doc, img)

    add_heading(doc, '3.5 Conclusion', level=2)
    doc.add_paragraph(d['conclusion'])
    doc.add_page_break()

    # 4. Naive Bayes
    d = REPORT_DATA['naive_bayes']
    add_heading(doc, f'4. {d["title"]}', level=1)
    doc.add_paragraph(f'Dataset: {d["dataset"]}')
    doc.add_paragraph(d['dataset_desc'])
    doc.add_paragraph(f'Algorithm: {d["algorithm"]}')
    doc.add_paragraph(f'Objective: {d["objective"]}')

    add_heading(doc, '4.1 Methodology', level=2)
    doc.add_paragraph(f'Features: {d["features"]}')
    doc.add_paragraph(f'Target: {d["target"]}')
    doc.add_paragraph(f'Preprocessing: {d["preprocessing"]}')
    doc.add_paragraph(
        'Text was converted to TF-IDF vectors with 5,000 features using unigrams and bigrams. '
        'A Multinomial Naive Bayes classifier with alpha=0.1 was trained on 24,000 samples and evaluated on 6,000.'
    )

    add_heading(doc, '4.2 Results', level=2)
    doc.add_paragraph(f'Test Accuracy: {d["test_accuracy"]}')

    add_heading(doc, '4.3 Top Shape Distribution', level=3)
    add_table(doc, ['Shape', 'Count'], d['top_shapes'])

    add_heading(doc, '4.4 Best Performing Classes', level=3)
    for cls, desc in d['best_classes']:
        p = doc.add_paragraph(style='List Bullet')
        run = p.add_run(f'{cls}: ')
        run.bold = True
        p.add_run(desc)

    add_heading(doc, '4.5 Visualizations', level=2)
    for img in d['images']:
        add_image(doc, img)

    add_heading(doc, '4.6 Conclusion', level=2)
    doc.add_paragraph(d['conclusion'])
    doc.add_page_break()

    # 5. Comparative Analysis
    add_heading(doc, '5. Comparative Analysis', level=1)
    add_table(doc,
        ['Metric', 'K-Means (LEGO)', 'KNN (Wildlife)', 'Naive Bayes (UFO)'],
        [
            ['Problem Type', 'Unsupervised', 'Classification', 'Text Classification'],
            ['Accuracy / Score', f'Silhouette: 0.52', f'89.8%', f'46.1%'],
            ['Best Parameters', 'k=4', 'k=8', 'alpha=0.1'],
            ['Key Challenge', 'Choosing optimal k', 'Class imbalance', 'Ambiguous text descriptions'],
        ]
    )

    add_heading(doc, '6. Conclusion', level=1)
    doc.add_paragraph(
        'Each algorithm demonstrated strengths suited to its respective dataset. K-Means effectively '
        'revealed natural groupings in LEGO sets. KNN achieved high accuracy on wildlife strike data '
        'but revealed class imbalance challenges. Naive Bayes proved capable of extracting signal from '
        'noisy text descriptions to predict UFO shapes. Together, they showcase the importance of '
        'matching algorithm choice to data characteristics and problem requirements.'
    )

    out = os.path.join(BASE, 'AI_Model_Training_Report.docx')
    doc.save(out)
    print(f'[Created] {out}')

# ═══════════════════════════════════════════
#  PPTX SLIDES
# ═══════════════════════════════════════════

def add_slide_title(prs, title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def add_slide_content(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
    return slide

def add_slide_image(prs, title, img_name, bullet=''):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    left = PInches(0.5); top = PInches(0.2); width = PInches(9); height = PInches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = PPt(28)
            r.font.bold = True
            r.font.color.rgb = PRGBColor(0x1A, 0x3C, 0x6E)
    # Image
    path = IMAGES.get(img_name)
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, PInches(0.8), PInches(1.2), width=PInches(8))
    # Bullet
    if bullet:
        txBox2 = slide.shapes.add_textbox(PInches(0.8), PInches(6.5), PInches(8), PInches(0.8))
        tf2 = txBox2.text_frame
        tf2.text = bullet
        for p in tf2.paragraphs:
            for r in p.runs:
                r.font.size = PPt(14)
                r.font.italic = True
    return slide

def build_slides():
    prs = Presentation()
    prs.slide_width = PInches(10)
    prs.slide_height = PInches(7.5)

    # Slide 1: Title
    s = add_slide_title(prs, 'AI Model Training Project',
                        'K-Means  ·  KNN  ·  Naive Bayes\n'
                        'LEGO Database  ·  Aviation Wildlife Strikes  ·  UFO Sightings')

    # Slide 2: Agenda
    add_slide_content(prs, 'Agenda', [
        '1. Project Overview',
        '2. Dataset Descriptions',
        '3. K-Means Clustering — LEGO Database',
        '4. K-Nearest Neighbors — Aviation Wildlife Strikes',
        '5. Naive Bayes — UFO Sightings Reports',
        '6. Comparative Analysis',
        '7. Conclusion & Future Work',
    ])

    # Slide 3: Overview
    add_slide_content(prs, 'Project Overview', [
        'Goal: Train and evaluate 3 ML algorithms on 3 distinct datasets',
        'Algorithms chosen to match problem characteristics:',
        '   • K-Means (unsupervised) → LEGO set clustering',
        '   • KNN (classification) → Wildlife strike risk',
        '   • Naive Bayes (text classification) → UFO shape prediction',
        'Tools: Python, scikit-learn, pandas, matplotlib',
    ])

    # Slide 4: Datasets
    add_slide_content(prs, 'Datasets Overview', [
        'LEGO Database (Rebrickable): 27K sets, year + piece count features',
        '   → Clustering by complexity and era',
        '',
        'Aviation Wildlife Strikes (FAA/OpenIntro): 19K records, 17 attributes',
        '   → Classifying strike effect (None / Precautions / Abort / Shutdown)',
        '',
        'UFO Sightings (NUFORC): 80K reports, text descriptions + shape labels',
        '   → Predicting UFO shape from witness descriptions',
    ])

    # Slide 5: K-Means Introduction
    add_slide_content(prs, 'K-Means Clustering — LEGO Database', [
        'Algorithm: Unsupervised K-Means with Euclidean distance',
        'Features: year (release year), num_parts (piece count)',
        'Preprocessing: StandardScaler normalization',
        '',
        'Optimal k determined via Elbow Method + Silhouette Score',
        '→ Optimal k = 4 (Silhouette Score: 0.52)',
        '',
        'Clusters discovered:',
        '  • Vintage sets (pre-1990, low parts)',
        '  • Modern small sets (accessories, polybags)',
        '  • Modern mid-size sets (500-1500 pieces)',
        '  • Large collector sets (2000+ pieces)',
    ])

    # Slide 6: K-Means Visuals
    add_slide_image(prs, 'K-Means: Elbow & Silhouette Analysis', 'lego_elbow.png',
                    'Elbow Method (left) and Silhouette Score (right) — k=4 is optimal')

    # Slide 7: K-Means Clusters
    add_slide_image(prs, 'K-Means: Cluster Visualization', 'lego_clusters.png',
                    'PCA projection (left) and Year vs Parts scatter (right) — 4 distinct groups visible')

    # Slide 8: KNN Introduction
    add_slide_content(prs, 'K-Nearest Neighbors — Wildlife Strikes', [
        'Algorithm: KNN Classification',
        'Features: height, speed, aircraft mass, engines, phase of flight,',
        '   time of day, sky condition, species (8 features)',
        'Target: Effect on flight (5 risk levels)',
        '',
        'Best k = 8 | Test Accuracy = 89.8% | CV Accuracy = 89.6%',
        '',
        'Challenge: Severe class imbalance (89.8% "None" class)',
        '  • "None" recall = 1.00 (excellent at identifying safe strikes)',
        '  • Rare classes (Engine Shut Down) need more data',
    ])

    # Slide 9: KNN Visuals
    add_slide_image(prs, 'KNN: Accuracy vs k and Confusion Matrix', 'wildlife_confusion_matrix.png',
                    'Left: accuracy plateaus around k=8. Right: confusion matrix shows majority class dominance.')

    # Slide 10: Naive Bayes Introduction
    add_slide_content(prs, 'Naive Bayes — UFO Sightings', [
        'Algorithm: Multinomial Naive Bayes + TF-IDF Vectorization',
        'Text preprocessing: lowercasing, punctuation/URL removal, stopwords',
        '5,000 TF-IDF features (unigrams + bigrams)',
        '10 shape classes to predict',
        '',
        'Test Accuracy = 46.1% (vs 24.3% majority class baseline)',
        '',
        'Best predicted shapes:',
        '  • Triangle (F1=0.61) — "silent", "triangular", "black"',
        '  • Fireball (F1=0.57) — "orange", "streak", "explosion"',
        '  • Light (F1=0.56) — most frequent class',
    ])

    # Slide 11: Naive Bayes Visuals
    add_slide_image(prs, 'Naive Bayes: Top Features per Shape', 'ufo_top_features.png',
                    'Most informative words for each UFO shape class')

    # Slide 12: Comparative
    add_slide_content(prs, 'Comparative Analysis', [
        '┌──────────────────┬────────────┬──────────┬──────────────┐',
        '│ Metric           │ K-Means    │ KNN      │ Naive Bayes  │',
        '├──────────────────┼────────────┼──────────┼──────────────┤',
        '│ Type             │ Unsuperv.  │ Classif. │ Text Classif.│',
        '│ Score            │ Silh. 0.52 │ 89.8%    │ 46.1%        │',
        '│ Best Param       │ k=4        │ k=8      │ alpha=0.1    │',
        '│ Key Challenge    │ Choosing k │ Imbalance│ Ambiguous txt│',
        '└──────────────────┴────────────┴──────────┴──────────────┘',
        '',
        'All algorithms successfully demonstrated core ML concepts',
        'Each revealed meaningful patterns in their respective datasets',
    ])

    # Slide 13: Conclusion
    add_slide_content(prs, 'Conclusion', [
        'K-Means effectively grouped LEGO sets by era and complexity',
        '  → Useful for market segmentation and trend analysis',
        '',
        'KNN achieved high accuracy on wildlife strike risk assessment',
        '  → Class imbalance remains a challenge for rare events',
        '',
        'Naive Bayes extracted meaningful signal from noisy text',
        '  → Text descriptions contain enough info to predict shape',
        '',
        'Future improvements:',
        '  • SMOTE / class weighting for imbalanced KNN',
        '  • Word embeddings (Word2Vec) for better text features',
        '  • Hierarchical clustering for deeper LEGO analysis',
    ])

    # Save
    out = os.path.join(BASE, 'AI_Model_Training_Slides.pptx')
    prs.save(out)
    print(f'[Created] {out}')

# ═══════════════════════════════════════════
#  RUN
# ═══════════════════════════════════════════

if __name__ == '__main__':
    build_report()
    build_slides()
    print('\nDone! Both files created.')
